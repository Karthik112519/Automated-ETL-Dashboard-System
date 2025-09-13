# report_gen.py - generate PPTX and simple PDF reports
import os, datetime
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import pandas as pd

def generate_pptx(cleaned_df, saved_plots_paths=None, out_path="report/ETL_Report.pptx"):
    saved_plots_paths = saved_plots_paths or []
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated ETL + Dashboard — Report"
    slide.placeholders[1].text = f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    # Data summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top rows"
    txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3))
    txt.text_frame.text = cleaned_df.head(8).to_string()
    # Stats slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Summary (numeric)"
    stats = cleaned_df.describe().round(3).to_string()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3))
    tx.text_frame.text = stats
    # Plots
    for p in saved_plots_paths:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = os.path.basename(p)
            slide.shapes.add_picture(p, Inches(1), Inches(1.5), width=Inches(8))
        except Exception:
            pass
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path

def generate_pdf(cleaned_df, out_path="report/ETL_Report.pdf"):
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    c = canvas.Canvas(out_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 72, "Automated ETL + Dashboard — Report")
    c.setFont("Helvetica", 10)
    c.drawString(72, height - 96, f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    # Top rows
    text = cleaned_df.head(10).to_string()
    y = height - 140
    lines = text.splitlines()
    c.setFont("Helvetica", 8)
    for line in lines:
        c.drawString(72, y, line)
        y -= 12
        if y < 72:
            c.showPage()
            y = height - 72
    c.showPage()
    c.save()
    return out_path